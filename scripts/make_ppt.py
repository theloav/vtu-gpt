from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

TITLE = "An AI-Driven Academic Assistant for Automated Calendar Extraction and Semantic Retrieval"
SUBTITLE = "Vel Tech Rangarajan Dr. Sagunthala R&D Institute of Science and Technology, Avadi, Chennai"
AUTHORS = "Authors: Mohib Rasool, Guru Prasad Reddy K, Dr. Rajendran T"

slides_content = [
    {
        "title": TITLE,
        "bullets": [
            SUBTITLE,
            AUTHORS,
            "Keywords: Academic Assistant, Event Extraction, Semantic Search, Calendar Automation",
        ],
        "notes": "Intro: project turns unstructured academic docs into answers and calendar events."
    },
    {
        "title": "Problem & Motivation",
        "bullets": [
            "Academic info scattered across PDFs, notices, emails, websites",
            "Students miss deadlines; staff answer repeat queries",
            "Manual calendar updates are slow and error-prone",
            "Need automation: extract events + answer questions instantly",
        ],
        "notes": "Pain points and why automation matters."
    },
    {
        "title": "System Overview",
        "bullets": [
            "Frontend: Chat, Calendar, Admin Upload",
            "Services: Auth, NLP/Chat, Document, Event, Sync",
            "AI: OpenAI (NLG), Pinecone (vector search), custom event extractor",
            "Data: SQLite (users/chats/events), Pinecone embeddings",
            "Integrations: Google Drive (5-min sync)",
        ],
        "notes": "Layered architecture overview."
    },
    {
        "title": "Methodology: Documents → Events",
        "bullets": [
            "Ingestion: PDF/DOCX/TXT; validation and text extraction",
            "Cleaning + semantic chunking",
            "Embeddings: 1536-dim, Pinecone storage",
            "Event extraction: dates, ranges, types, confidence",
            "Calendar update: title, date, type, source",
        ],
        "notes": "Core pipeline."
    },
    {
        "title": "Implementation Highlights",
        "bullets": [
            "Stack: React/Next.js, Node/Express, SQLite, Pinecone, OpenAI",
            "Extractor: flexible regex + keywords + confidence",
            "Sync: Google Drive, batching, retries",
            "Security: JWT, validation, parameterized SQL",
            "Admin tools: Upload dashboard, events API",
        ],
        "notes": "Practical build aspects."
    },
    {
        "title": "Evaluation & Results",
        "bullets": [
            "Dataset: calendars, timetables, announcements",
            "Event extraction: Precision 94.2%, Recall 91.7%, F1 92.9%",
            "Performance: p95 < 2.8s; stable up to 200 users",
            "User study: 4.6/5 satisfaction; 67% faster info finding",
            "Uptime: 99.7% during pilot",
        ],
        "notes": "Impact metrics."
    },
    {
        "title": "Demo / Screenshots",
        "bullets": [
            "Chat UI answering grounded queries",
            "Calendar view with extracted events",
            "Admin upload and processing status",
            "(Add 2–3 screenshots here)",
        ],
        "notes": "Walk through screenshots briefly."
    },
    {
        "title": "Conclusion & Next Steps",
        "bullets": [
            "Automated event extraction + semantic search + conversational access",
            "Impact: timely info, reduced workload, higher satisfaction",
            "Next: multi-language, advanced event reasoning, LMS/SIS, mobile apps",
            "Thank you! Questions?",
        ],
        "notes": "Close and invite questions."
    },
]

def make_presentation(path="slides/Project-Presentation.pptx"):
    prs = Presentation()

    # Title slide layout (0)
    title_layout = prs.slide_layouts[0]
    # Title and Content layout (1)
    content_layout = prs.slide_layouts[1]

    for idx, sc in enumerate(slides_content):
        if idx == 0:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = sc["title"]
            subtitle = slide.placeholders[1]
            subtitle.text = "\n".join(sc["bullets"])  # stacked lines
        else:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = sc["title"]
            body = slide.placeholders[1].text_frame
            body.clear()
            for i, b in enumerate(sc["bullets"]):
                if i == 0:
                    p = body.paragraphs[0]
                else:
                    p = body.add_paragraph()
                p.text = b
                p.level = 0
            
        # Add speaker notes
        notes = slide.notes_slide.notes_text_frame
        notes.text = sc["notes"]

    prs.save(path)
    return path

if __name__ == "__main__":
    out = make_presentation()
    print(f"Saved PPTX to {out}")

